import re
import copy
import zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Run from the repo root:
#   python3 signedby-app/scripts/build-deck-slides-business.py
#
# Adds ONE new "BUSINESS" slide (direct ask, 2026-08-08: "Make an extra
# slide for the pitch deck, 'Designed for business' ... business customer
# branding, payment link and gated access features and the pre-signing and
# post-signing brand and ad placement opportunity") to BOTH the pitch and
# product decks, inserted right after DEVELOPERS in each — same "more
# platform depth" cluster (PRODUCT/SIGN/SEAL/VERIFY/DEVELOPERS) rather than
# the business-case slides that follow it in the pitch deck (PROOF/MARKET/
# WHY NOW/etc).
#
# Baseline files (direct instruction, same day): v12 of both decks, which
# were restructured outside this script (SEAL/VERIFY consolidated a Verify
# slide differently than build-deck-slides-v9.py originally shipped it) --
# confirmed the SEAL slide's shape ids {2,3,4,19,20} = eyebrow/title/
# subtitle/footer/pagenum still hold in both v12 files before relying on
# them here.
#
# The three feature bullets (branding, payment collection, gated access)
# are all genuinely Business-plan-only today (src/lib/plan.ts: branding,
# customBranding, paymentCollection, docGate all gate to ["business"]) --
# not a stretch to frame this as "designed for business". The ad/brand
# placement angle is explicitly NOT something that's built -- framed as an
# idea/opportunity, visually lighter-weight than the three real bullets, so
# the slide doesn't overclaim a monetization feature that doesn't exist.

TARGETS = [
    ("SignedBy_Pitch_Teaser_v12.pptx", "SignedBy_Pitch_Teaser_v13.pptx"),
    ("SignedBy_Product_Teaser_v12.pptx", "SignedBy_Product_Teaser_v13.pptx"),
]
APP_PUBLIC = "signedby-app/public"

SLIDE_W = 12192000
MARGIN_L = 548640
MARGIN_R = 548640
FOOTER_T = 6446520
TOP = 2346350

NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

_next_slide_n = [None]


def max_slide_number_in_file(path):
    """See build-deck-slides-v9.py's docstring for why this reads the raw
    zip XML directly instead of asking python-pptx's object model -- same
    corruption footgun, same fix, reused verbatim here."""
    with zipfile.ZipFile(path) as z:
        try:
            data = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        except KeyError:
            return 0
    nums = [int(n) for n in re.findall(r'slide(\d+)\.xml"', data)]
    return max(nums) if nums else 0


def dedupe_new_slide_partname(new_slide):
    from pptx.opc.packuri import PackURI
    _next_slide_n[0] += 1
    new_slide.part.partname = PackURI(f"/ppt/slides/slide{_next_slide_n[0]}.xml")


def duplicate_bg_and_ids(prs, source_slide, keep_ids):
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    dedupe_new_slide_partname(new_slide)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    source_bg = source_slide._element.find(f".//{NS_P}cSld/{NS_P}bg")
    if source_bg is not None:
        cSld = new_slide._element.find(f"{NS_P}cSld")
        spTree = cSld.find(f"{NS_P}spTree")
        cSld.insert(list(cSld).index(spTree), copy.deepcopy(source_bg))

    for shp in source_slide.shapes:
        if shp.shape_id in keep_ids:
            new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def set_text(slide, shape_id, text):
    for shp in slide.shapes:
        if shp.shape_id == shape_id and shp.has_text_frame:
            para = shp.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra._r.getparent().remove(extra._r)
            else:
                run = para.add_run()
                run.text = text
            return
    raise ValueError(f"shape_id {shape_id} not found")


def add_text(slide, l, t, w, h, text, size_pt, color_rgb, bold=False, italic=False, align="l", font="Calibri"):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "c":
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color_rgb
    return box


def add_caption(slide, l, t, w, text):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(280000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.name = "Calibri"
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def fit_in_box(img_path, max_w, max_h):
    with Image.open(img_path) as im:
        w, h = im.size
    scale = min(max_w / w, max_h / h)
    return round(w * scale), round(h * scale)


def add_feature_bullet(slide, row_t, icon_path, title, body, icon_l=8001900):
    """SEAL slide's exact icon-badge geometry (white circle + inset icon +
    title/body to the right) -- reused as-is rather than reinvented, since
    it's already the deck's established "full-weight feature bullet"
    pattern (as opposed to VERIFY's smaller plain-checkmark takeaways)."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(icon_l), Emu(row_t), Emu(502920), Emu(502920))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    circle.line.fill.background()
    icon_inset = (502920 - 384048) // 2
    slide.shapes.add_picture(icon_path, Emu(icon_l + icon_inset), Emu(row_t + icon_inset), Emu(384048), Emu(384048))
    text_l = icon_l + 685800
    text_w = SLIDE_W - MARGIN_R - text_l
    add_text(slide, text_l, row_t - 30000, text_w, 365760, title, 15.5, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    add_text(slide, text_l, row_t + 228600, text_w, 1234440, body, 12, RGBColor(0x94, 0xA3, 0xB8))


def build_business_slide(prs, template_slide):
    keep_ids = {2, 3, 4, 19, 20}
    slide = duplicate_bg_and_ids(prs, template_slide, keep_ids)
    set_text(slide, 2, "BUSINESS")
    set_text(slide, 3, "Every business document is also a brand moment.")
    set_text(slide, 4, "Branding, payment collection, and gated next-step access ship today, Business plan only "
                       "— the same two screens are also the moments a signer is most engaged.")

    # 3-column layout (direct ask, 2026-08-08, revised from an earlier
    # 2-image-plus-callout draft): pre-signing screenshot | post-signing
    # screenshot | 3 feature bullets. Both hero images are REAL product
    # screenshots the user captured on their own phone and sent directly
    # (hero-business-presign-real.png / hero-business-postsign-real.png),
    # not the next/og mockups this script originally generated
    # (hero-business-presign-header.png / -postsign-confirmation.png,
    # still in scripts/ but unused now that real screenshots exist —
    # matches this deck's standing preference for real product captures
    # over recreated mockups wherever one is available). The pre-signing
    # shot had the user's own name ("Signing as Michael Eagles") redacted
    # with a solid bar before cropping (direct instruction); the
    # post-signing shot is a different real screenshot the user supplied
    # specifically because it doesn't show the payment/certificate/QR
    # block, for a cleaner image ("without the QR code so it looks
    # cleaner") — it shows a different demo signer name (Amara Okafor,
    # not the user), so no redaction was needed there.
    content_w = SLIDE_W - MARGIN_L - MARGIN_R
    gap = 274320
    col_w = (content_w - 2 * gap) // 3
    col1_l = MARGIN_L
    col2_l = col1_l + col_w + gap
    col3_l = col2_l + col_w + gap
    max_h = FOOTER_T - TOP

    presign_path = f"{APP_PUBLIC}/hero-business-presign-real.png"
    presign_w, presign_h = fit_in_box(presign_path, col_w, max_h)
    slide.shapes.add_picture(presign_path, Emu(col1_l), Emu(TOP), Emu(presign_w), Emu(presign_h))
    add_caption(slide, col1_l, TOP + presign_h + 40000, presign_w,
                "The real pre-signing screen — a live product screenshot, not a mockup.")

    postsign_path = f"{APP_PUBLIC}/hero-business-postsign-real.png"
    postsign_w, postsign_h = fit_in_box(postsign_path, col_w, max_h)
    slide.shapes.add_picture(postsign_path, Emu(col2_l), Emu(TOP), Emu(postsign_w), Emu(postsign_h))
    add_caption(slide, col2_l, TOP + postsign_h + 40000, postsign_w,
                "The real post-signing confirmation — a live product screenshot, not a mockup.")

    # 3 real feature bullets -- exact SEAL-slide geometry/spacing, aligned
    # to this slide's own 3rd column instead of SEAL's hardcoded L8.75in.
    features = [
        (f"{APP_PUBLIC}/deck-icons/Palette.png", "Full white-label branding",
         "Logo, brand color, and org name across the signing page and confirmation screen — SignedBy's own mark never shows."),
        (f"{APP_PUBLIC}/deck-icons/CreditCard.png", "Payment collection, built in",
         "A Pay Now prompt lands right on the confirmation screen — no separate invoicing step or link to send."),
        (f"{APP_PUBLIC}/deck-icons/ShieldCheck.png", "Gated next-step access",
         "A next-step link — onboarding portal, dataroom, download — only unlocks once every signer has completed the document."),
    ]
    for (icon, title, body), row_t in zip(features, (2286000, 3703320, 5120640)):
        add_feature_bullet(slide, row_t, icon, title, body, icon_l=col3_l)

    return slide


def move_slide_after(prs, slide, after_index):
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    anchor = sld_ids[after_index]
    target = None
    for sldId in xml_slides:
        if int(sldId.get("id")) == slide.slide_id:
            target = sldId
            break
    xml_slides.remove(target)
    anchor_pos = list(xml_slides).index(anchor)
    xml_slides.insert(anchor_pos + 1, target)


def build_for(src, dst):
    _next_slide_n[0] = max_slide_number_in_file(src)
    prs = Presentation(src)

    template_slide = None
    developers_slide = None
    for s in prs.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t == "SEAL" and template_slide is None:
                    template_slide = s
                if t == "DEVELOPERS" and developers_slide is None:
                    developers_slide = s
                break

    business_slide = build_business_slide(prs, template_slide)
    dev_index = list(prs.slides).index(developers_slide)
    move_slide_after(prs, business_slide, after_index=dev_index)

    for idx, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit():
                shp.text_frame.paragraphs[0].runs[0].text = str(idx)

    prs.save(dst)
    print("saved", dst, "slides:", len(prs.slides._sldIdLst))


def main():
    for src, dst in TARGETS:
        build_for(src, dst)


if __name__ == "__main__":
    main()
