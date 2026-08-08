import copy
import re
import zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

# Run from the repo root:
#   python3 signedby-app/scripts/build-deck-slides-v9.py
#
# Adds 2 new slides (direct ask, 2026-08-08), inserted after SEAL and
# before DEVELOPERS:
#   - CAPABILITIES: the real Sent/Sealed confirmation popovers from
#     send-seal-transition.tsx side by side, plus 4 core-capability blurbs.
#   - VERIFICATION: the real standalone Verified Badge asset next to a
#     mockup of the /verify page's green result state, plus 3 takeaways.

SRC = "SignedBy_Product_Teaser_v8.pptx"
DST = "SignedBy_Product_Teaser_v9.pptx"
APP_PUBLIC = "signedby-app/public"

SLIDE_W = 12192000
MARGIN_L = 548640
MARGIN_R = 548640
FOOTER_T = 6446520
CAPTION_GAP = 340000
TOP = 2346350

NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


_next_slide_n = [None]  # set by main() from max_slide_number_in_file(SRC); a
                         # 1-item list so nested functions can mutate it
                         # without a `global` statement.


def max_slide_number_in_file(path):
    """Ground-truth scan of the SOURCE file's raw presentation.xml.rels XML
    for the highest existing /ppt/slides/slideN.xml target number.

    Earlier versions of this guard asked python-pptx itself (first via
    package.iter_parts(), then via prs.part.rels) which slide partnames
    were already in use, and picked one number higher. Both turned out to
    be unreliable: empirically, simply iterating `for s in prs.slides: for
    shp in s.shapes: ...` (e.g. to look up a slide by its eyebrow text --
    exactly what this script's own main() does) was enough to make
    prs.part.rels report the SAME relationship set with 3 partnames
    duplicated, on an otherwise byte-identical reload of the same file --
    a python-pptx-internal lazy-loading/caching quirk, not anything this
    script's own slide-building does. Reading the raw XML straight from
    the zip, once, before python-pptx's object model ever touches the
    relationship graph, sidesteps that entirely."""
    with zipfile.ZipFile(path) as z:
        try:
            data = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        except KeyError:
            return 0
    nums = [int(n) for n in re.findall(r'slide(\d+)\.xml"', data)]
    return max(nums) if nums else 0


def duplicate_bg_and_ids(prs, source_slide, keep_ids):
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    dedupe_new_slide_partname(new_slide)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    source_bg = source_slide._element.find(f".//{NS_P}cSld/{NS_P}bg")
    if source_bg is not None:
        cSld = new_slide._element.find(f"{NS_P}cSld")
        spTree = cSld.find(f"{NS_P}spTree")
        cSld.insert(list(cSld).index(spTree), copy.deepcopy(source_bg))

    for shp in source_slide.shapes:
        if shp.shape_id in keep_ids:
            new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def set_text(slide, shape_id, text):
    for shp in slide.shapes:
        if shp.shape_id == shape_id and shp.has_text_frame:
            para = shp.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra._r.getparent().remove(extra._r)
            else:
                run = para.add_run()
                run.text = text
            return
    raise ValueError(f"shape_id {shape_id} not found")


def add_text(slide, l, t, w, h, text, size_pt, color_rgb, bold=False, italic=False, align="l", font="Calibri"):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "c":
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color_rgb
    return box


def fit_in_box(img_path, max_w, max_h):
    with Image.open(img_path) as im:
        w, h = im.size
    scale = min(max_w / w, max_h / h)
    return round(w * scale), round(h * scale)


def add_caption(slide, l, t, w, text):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(280000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.name = "Calibri"
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


# ---------------------------------------------------------------------------
# CAPABILITIES slide
# ---------------------------------------------------------------------------

def build_capabilities_slide(prs, template_slide):
    keep_ids = {2, 3, 4, 19, 20}
    slide = duplicate_bg_and_ids(prs, template_slide, keep_ids)
    set_text(slide, 2, "CAPABILITIES")
    set_text(slide, 3, "Every document ends one of two ways: sent, or sealed.")
    set_text(slide, 4, "The same two confirmations customers see in the product today, backed by the capabilities underneath.")

    # Hero: the real Sent/Sealed popovers, side by side, centered.
    hero_path = f"{APP_PUBLIC}/hero-send-seal-popovers.png"
    max_w, max_h = 7800000, 2500000
    pic_w, pic_h = fit_in_box(hero_path, max_w, max_h)
    pic_l = (SLIDE_W - pic_w) // 2
    pic_t = TOP
    slide.shapes.add_picture(hero_path, Emu(pic_l), Emu(pic_t), Emu(pic_w), Emu(pic_h))
    add_caption(slide, pic_l, pic_t + pic_h + 40000, pic_w,
                "The real Sent / Sealed confirmation popovers, as shown in the product today — not a mockup.")

    # 4 core-capability blurbs underneath, same 4-column rhythm as the
    # USE CASES slide (build-deck-slides-v8.py) but text-only (no badge
    # icon -- the popovers above already carry the icon weight).
    n_cols = 4
    col_gap = 300000
    col_w = (SLIDE_W - MARGIN_L - MARGIN_R - (n_cols - 1) * col_gap) // n_cols
    row_t = pic_t + pic_h + 420000
    title_h = 340000
    body_h = 1300000

    capabilities = [
        ("Legally binding e-signatures",
         "An eIDAS-compliant signing flow built for EU businesses and their customers."),
        ("Cryptographic sealing",
         "A SHA-512 file hash plus an RFC 3161 trusted timestamp anchor every sealed document."),
        ("Per-signer verification",
         "Email OTP authentication confirms who actually opened and signed, not just a link click."),
        ("Public verification",
         "Anyone can check a sealed document is genuine — no SignedBy account required."),
    ]

    for i, (title, body) in enumerate(capabilities):
        col_l = MARGIN_L + i * (col_w + col_gap)
        add_text(slide, col_l, row_t, col_w, title_h, title, 15, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        add_text(slide, col_l, row_t + title_h + 60000, col_w, body_h, body, 11.5, RGBColor(0x94, 0xA3, 0xB8))

    return slide


# ---------------------------------------------------------------------------
# VERIFICATION slide
# ---------------------------------------------------------------------------

def build_verification_slide(prs, template_slide):
    keep_ids = {2, 3, 4, 19, 20}
    slide = duplicate_bg_and_ids(prs, template_slide, keep_ids)
    set_text(slide, 2, "VERIFICATION")
    set_text(slide, 3, "Anyone can check a sealed document is real — no account needed.")
    set_text(slide, 4, "Scan the badge or paste the checksum; the verify page confirms it in seconds.")

    max_h = 3550000

    badge_path = f"{APP_PUBLIC}/hero-verified-badge.png"
    verify_path = f"{APP_PUBLIC}/hero-verify-result.png"

    badge_w, badge_h = fit_in_box(badge_path, 2700000, max_h)
    verify_w, verify_h = fit_in_box(verify_path, 3600000, max_h)

    gap = 300000
    badge_l = MARGIN_L
    verify_l = badge_l + badge_w + gap
    img_t = TOP

    slide.shapes.add_picture(badge_path, Emu(badge_l), Emu(img_t), Emu(badge_w), Emu(badge_h))
    add_caption(slide, badge_l, img_t + badge_h + 40000, badge_w,
                "The real Verified Badge asset — not a mockup.")

    slide.shapes.add_picture(verify_path, Emu(verify_l), Emu(img_t), Emu(verify_w), Emu(verify_h))
    add_caption(slide, verify_l, img_t + verify_h + 40000, verify_w,
                "Illustrative mockup of the verify page's result — not a screenshot.")

    # 3 takeaways in the remaining right column, small green-check bullet
    # to match the verification theme.
    takeaways_l = verify_l + verify_w + gap
    takeaways_w = SLIDE_W - MARGIN_R - takeaways_l
    check_wh = 260000
    row_gap = 1080000
    row_t0 = img_t + 60000

    takeaways = [
        ("No account needed",
         "Scan the badge or paste the checksum — anyone can check a document is genuine."),
        ("Confirms two separate facts",
         "That the file is unaltered, and that whoever sealed it passed identity verification."),
        ("Backed by real cryptography",
         "A SHA-512 hash and an RFC 3161 trusted timestamp — not just a claim in our database."),
    ]

    for i, (title, body) in enumerate(takeaways):
        row_t = row_t0 + i * row_gap
        slide.shapes.add_picture(f"{APP_PUBLIC}/deck-icons/Check-emerald.png", Emu(takeaways_l), Emu(row_t), Emu(check_wh), Emu(check_wh))
        add_text(slide, takeaways_l + check_wh + 140000, row_t - 40000, takeaways_w - check_wh - 140000, 400000,
                 title, 15, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        add_text(slide, takeaways_l + check_wh + 140000, row_t + 300000, takeaways_w - check_wh - 140000, 700000,
                 body, 11.5, RGBColor(0x94, 0xA3, 0xB8))

    return slide


def move_slides_after(prs, slides_to_move, after_index):
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    anchor = sld_ids[after_index]
    target_sldIds = []
    for slide in slides_to_move:
        for sldId in xml_slides:
            if int(sldId.get("id")) == slide.slide_id:
                target_sldIds.append(sldId)
                break
    for sldId in target_sldIds:
        xml_slides.remove(sldId)
    anchor_pos = list(xml_slides).index(anchor)
    for offset, sldId in enumerate(target_sldIds):
        xml_slides.insert(anchor_pos + 1 + offset, sldId)


def dedupe_new_slide_partname(new_slide):
    """add_slide() names new slide parts via a naive len(sldIdLst)+1 scheme
    (pptx/parts/presentation.py's _next_slide_partname) that collides with
    an existing (different) part once the deck has any non-contiguous
    partname history (e.g. slides removed+added across earlier rebuilds).

    Two earlier attempts at guarding this both asked python-pptx's own
    object model which slide partnames were already in use -- first
    package.iter_parts(), then prs.part.rels -- and both turned out to be
    unreliable in practice (see max_slide_number_in_file's docstring for
    the specifics). This version just assigns the next number from a
    plain counter (_next_slide_n) seeded once, up front, from a raw XML
    scan of the source file -- no python-pptx object-graph involved at
    all, so there's nothing left for it to get inconsistent about."""
    from pptx.opc.packuri import PackURI

    _next_slide_n[0] += 1
    new_slide.part.partname = PackURI(f"/ppt/slides/slide{_next_slide_n[0]}.xml")


def main():
    _next_slide_n[0] = max_slide_number_in_file(SRC)
    prs = Presentation(SRC)
    slides_by_eyebrow = {}
    for s in prs.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t in ("PRODUCT", "SEAL"):
                    slides_by_eyebrow.setdefault(t, s)
                    break

    template_slide = slides_by_eyebrow["PRODUCT"]
    seal_slide = slides_by_eyebrow["SEAL"]

    capabilities_slide = build_capabilities_slide(prs, template_slide)
    verification_slide = build_verification_slide(prs, template_slide)

    seal_index = list(prs.slides).index(seal_slide)
    move_slides_after(prs, [capabilities_slide, verification_slide], after_index=seal_index)

    slides_in_order = list(prs.slides)
    for idx, slide in enumerate(slides_in_order, start=1):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit():
                shp.text_frame.paragraphs[0].runs[0].text = str(idx)

    prs.save(DST)
    print("saved", DST, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
