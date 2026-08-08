import copy
from pptx import Presentation
from pptx.util import Emu
from PIL import Image

SRC = "SignedBy_Product_Teaser_v6.pptx"
DST = "SignedBy_Product_Teaser_v7.pptx"
APP_PUBLIC = "signedby-app/public"

NAVY = "0F172A"

# ---- geometry constants lifted from slide 5 (PRODUCT), the template slide ----
CARD_L = 624112
CARD_T = 2346350
CARD_PAD = 120000
FOOTER_T = 6446520
CAPTION_GAP = 340000  # room for an honesty caption between hero card and footer
MAX_CARD_W = 6656369  # same width slide 5's own card uses
MAX_CARD_H = FOOTER_T - CAPTION_GAP - CARD_T

ICON_ROWS = [2286000, 3703320, 5120640]  # T of each circle/icon/title/body row
ICON_CIRCLE_L = 8001000
ICON_PIC_L = 8061350
ICON_PIC_OFFSET_T = 60350  # icon pic T - row T, constant across rows (2346350-2286000)
TEXT_L = 8686800
TEXT_W = 2971800
TITLE_OFFSET_T = -27432  # title T - row T (2258568 - 2286000)
BODY_OFFSET_T = 204430   # body T - row T (2490430 - 2286000)
TITLE_H = 365760
BODY_H = 1234440
ICON_WH = 382219
CIRCLE_WH = 502920


def duplicate_template_slide(prs, source_slide):
    """New slide on the same layout, with all of source_slide's shapes
    except its PICTUREs copied over (positions/fonts preserved), pictures
    re-added by the caller since each new slide uses different images."""
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # Slide-level background (navy fill) lives as <p:bg> on <p:cSld>, a
    # sibling of <p:spTree> -- NOT one of slide.shapes, so the shape-copy
    # loop below never touches it. Without this, new slides silently fall
    # back to the layout's white bg1 background. Copy it explicitly and
    # insert before spTree (schema requires bg to precede spTree).
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    source_bg = source_slide._element.find(f".//{ns}cSld/{ns}bg")
    if source_bg is not None:
        cSld = new_slide._element.find(f"{ns}cSld")
        spTree = cSld.find(f"{ns}spTree")
        cSld.insert(list(cSld).index(spTree), copy.deepcopy(source_bg))

    id_map = {}
    for shp in source_slide.shapes:
        if shp.shape_type == 13:  # PICTURE — caller re-adds these
            continue
        new_el = copy.deepcopy(shp._element)
        new_slide.shapes._spTree.append(new_el)
        id_map[shp.shape_id] = new_el
    return new_slide


def set_text(slide, shape_id, text):
    for shp in slide.shapes:
        if shp.shape_id == shape_id and shp.has_text_frame:
            para = shp.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra._r.getparent().remove(extra._r)
            else:
                run = para.add_run()
                run.text = text
            return
    raise ValueError(f"shape_id {shape_id} not found")


def fit_hero(img_path):
    with Image.open(img_path) as im:
        w, h = im.size
    pic_w_avail = MAX_CARD_W - 2 * CARD_PAD
    pic_h_avail = MAX_CARD_H - 2 * CARD_PAD
    scale = min(pic_w_avail / w, pic_h_avail / h)
    pic_w, pic_h = round(w * scale), round(h * scale)
    card_w, card_h = pic_w + 2 * CARD_PAD, pic_h + 2 * CARD_PAD
    card_t = CARD_T + round((MAX_CARD_H - card_h) / 2)
    return {
        "card": (CARD_L, card_t, card_w, card_h),
        "pic": (CARD_L + CARD_PAD, card_t + CARD_PAD, pic_w, pic_h),
    }


def set_xfrm(shape_el, l, t, w, h):
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    xfrm = shape_el.find(f".//{ns}xfrm")
    off = xfrm.find(f"{ns}off")
    ext = xfrm.find(f"{ns}ext")
    off.set("x", str(l))
    off.set("y", str(t))
    ext.set("cx", str(w))
    ext.set("cy", str(h))


def add_hero(slide, source_card_shape, img_path):
    geo = fit_hero(img_path)
    card_el = copy.deepcopy(source_card_shape._element)
    slide.shapes._spTree.append(card_el)
    l, t, w, h = geo["card"]
    set_xfrm(card_el, l, t, w, h)
    l, t, w, h = geo["pic"]
    slide.shapes.add_picture(img_path, Emu(l), Emu(t), Emu(w), Emu(h))
    return geo["card"]


def add_caption(slide, card_geo, text):
    from pptx.util import Pt
    l, t, w, h = card_geo
    box = slide.shapes.add_textbox(Emu(l), Emu(t + h + 60000), Emu(w), Emu(260000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.name = "Calibri"
    run.font.italic = True
    from pptx.dml.color import RGBColor
    run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def build_feature_slide(prs, template_slide, source_card_shape, source_icon_pics, eyebrow, title, intro,
                         hero_path, caption, callouts):
    """callouts: list of (icon_png_path, title, body) length 3."""
    slide = duplicate_template_slide(prs, template_slide)
    set_text(slide, 2, eyebrow)
    set_text(slide, 3, title)
    set_text(slide, 4, intro)

    title_ids = [9, 13, 17]
    body_ids = [10, 14, 18]
    for i, (icon_path, ctitle, cbody) in enumerate(callouts):
        set_text(slide, title_ids[i], ctitle)
        set_text(slide, body_ids[i], cbody)
        row_t = ICON_ROWS[i]
        slide.shapes.add_picture(icon_path, Emu(ICON_PIC_L), Emu(row_t + ICON_PIC_OFFSET_T), Emu(ICON_WH), Emu(ICON_WH))

    card_geo = add_hero(slide, source_card_shape, hero_path)
    add_caption(slide, card_geo, caption)
    # footer text (shape 19) stays "SignedBy · Pre-Seed" already copied; page number (20) set later in renumber pass
    return slide


def move_slides_after(prs, slides_to_move, after_index):
    """after_index is 0-based index (in the CURRENT xml order, before moving)
    of the slide that the moved slides should immediately follow."""
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    anchor = sld_ids[after_index]
    for slide in slides_to_move:
        el = slide._element.getparent()  # not used; slide id element lookup below
    # Build map from slide part to sldId element
    id_by_rId = {}
    for sldId in sld_ids:
        id_by_rId[sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")] = sldId
    for slide in slides_to_move:
        rId = prs.part.relate_to(slide.part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide", reuse=True) if False else None
    # Simpler: locate each new slide's sldId by matching slide.slide_id
    target_sldIds = []
    for slide in slides_to_move:
        for sldId in xml_slides:
            if int(sldId.get("id")) == slide.slide_id:
                target_sldIds.append(sldId)
                break
    for sldId in target_sldIds:
        xml_slides.remove(sldId)
    anchor_pos = list(xml_slides).index(anchor)
    for offset, sldId in enumerate(target_sldIds):
        xml_slides.insert(anchor_pos + 1 + offset, sldId)


def main():
    prs = Presentation(SRC)
    template_slide = prs.slides[4]  # PRODUCT (0-indexed slide 5)
    source_card_shape = None
    icon_pics = []
    for shp in template_slide.shapes:
        if shp.shape_id == 5:
            source_card_shape = shp
    assert source_card_shape is not None

    sign_slide = build_feature_slide(
        prs, template_slide, source_card_shape, icon_pics,
        eyebrow="SIGN",
        title="Rental signing that fits in the time it takes to hand over the keys.",
        intro="No app for the renter to download — scan a QR code, sign on their own phone, and drive off in minutes.",
        hero_path=f"{APP_PUBLIC}/hero-boat-jet-ski-rental-qr-signing.png",
        caption="Illustrative mockup of the shipped QR-to-sign flow — not a screenshot.",
        callouts=[
            (f"{APP_PUBLIC}/deck-icons/QrCode.png", "Scan to sign, no app",
             "The renter scans a QR code with their phone camera and signs directly in the browser — nothing to install."),
            (f"{APP_PUBLIC}/deck-icons/Clock.png", "Signed before they leave the counter",
             "The whole flow takes under a minute end to end, so operators aren't chasing signatures after the fact."),
            (f"{APP_PUBLIC}/deck-icons/ListChecks.png", "Every field already in place",
             "Fixed-fee and per-day rental terms are mapped to the right fields ahead of time — nothing to configure per rental."),
        ],
    )

    seal_slide = build_feature_slide(
        prs, template_slide, source_card_shape, icon_pics,
        eyebrow="SEAL",
        title="Proof an invoice hasn't been touched since it was sent.",
        intro="A verification badge — hash, trusted timestamp, and a public verify page — embedded right on the document.",
        hero_path=f"{APP_PUBLIC}/hero-verified-badge-invoice.png",
        caption="Illustrative mockup of the shipped Seal / verify flow — not a screenshot.",
        callouts=[
            (f"{APP_PUBLIC}/deck-icons/Stamp.png", "Sealed, not just signed",
             "A cryptographic hash and an RFC 3161 trusted timestamp anchor the exact bytes of the invoice the moment it's sealed."),
            (f"{APP_PUBLIC}/deck-icons/ScanLine.png", "Anyone can verify it",
             "Scanning the badge opens a public page confirming the document is unchanged since it was sealed — no account needed."),
            (f"{APP_PUBLIC}/deck-icons/ShieldAlert.png", "Built to stop invoice fraud",
             "Protects freelancers and agencies against the classic scam: a client quietly editing bank details on a copy of a real invoice."),
        ],
    )

    draft_slide = build_feature_slide(
        prs, template_slide, source_card_shape, icon_pics,
        eyebrow="DRAFT",
        title="A working contract from a plain-English description.",
        intro="Describe the deal, get a structured starting draft, then edit it like any other SignedBy document.",
        hero_path=f"{APP_PUBLIC}/hero-ai-draft-mockup.png",
        caption="Illustrative mockup of the AI Draft flow — not a screenshot.",
        callouts=[
            (f"{APP_PUBLIC}/deck-icons/WandSparkles.png", "Describe it, get a draft",
             "A plain-English description in, a structured agreement out — no template library to search first."),
            (f"{APP_PUBLIC}/deck-icons/FileText.png", "Real terms, not filler",
             "Drafts include the specific terms mentioned — payment, ownership, deadlines — not generic boilerplate."),
            (f"{APP_PUBLIC}/deck-icons/PenLine.png", "Then edit like any document",
             "The draft drops straight into the same field editor and signer flow as everything else on SignedBy."),
        ],
    )

    # Reorder: new slides go right after PRODUCT (original 0-index 4)
    move_slides_after(prs, [sign_slide, seal_slide, draft_slide], after_index=4)

    # ---- fix known pre-existing bugs while we're in here ----
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame and "Oonk—" in shp.text_frame.text:
                shp.text_frame.paragraphs[0].runs[0].text = shp.text_frame.paragraphs[0].runs[0].text.replace(
                    "Oonk—", "Oonk —"
                )

    # ---- renumber every footer page number sequentially ----
    slides_in_order = list(prs.slides)
    for idx, slide in enumerate(slides_in_order, start=1):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit():
                shp.text_frame.paragraphs[0].runs[0].text = str(idx)

    prs.save(DST)
    print("saved", DST, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
