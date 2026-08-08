import copy
import re
import zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# Run from the repo root (same convention as build-deck-slides-v7.py):
#   python3 signedby-app/scripts/build-deck-slides-v8.py
#
# Replaces the standalone SIGN / SEAL / DRAFT feature slides added in v7
# with a leaner 3-slide sequence (direct ask, 2026-08-08): one 4-up
# SIGN/SEAL/QUOTE/DRAFT use-case summary slide (headers reuse the real
# app's own New Document tab badges), then a zoom slide each for just
# Sign and Seal (kept, same content as v7's). Draft keeps its one line in
# the summary row but loses its dedicated zoom slide -- Quote is net-new,
# was not covered anywhere in the deck before this.

SRC = "SignedBy_Product_Teaser_v7.pptx"
DST = "SignedBy_Product_Teaser_v8.pptx"
APP_PUBLIC = "signedby-app/public"

# ---- geometry constants lifted from slide 5 (PRODUCT), the template slide ----
CARD_L = 624112
CARD_T = 2346350
CARD_PAD = 120000
FOOTER_T = 6446520
CAPTION_GAP = 340000
MAX_CARD_W = 6656369
MAX_CARD_H = FOOTER_T - CAPTION_GAP - CARD_T

ICON_ROWS = [2286000, 3703320, 5120640]
ICON_PIC_L = 8061350
ICON_PIC_OFFSET_T = 60350
TEXT_L = 8686800
ICON_WH = 382219

NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


_next_slide_n = [None]  # set by main() from max_slide_number_in_file(SRC); a
                         # 1-item list so nested functions can mutate it
                         # without a `global` statement.


def max_slide_number_in_file(path):
    """Ground-truth scan of the SOURCE file's raw presentation.xml.rels XML
    for the highest existing /ppt/slides/slideN.xml target number.

    UPDATED 2026-08-08 (again): this function used to ask python-pptx's
    own object model instead -- first package.iter_parts(), then
    prs.part.rels -- and BOTH turned out to be unreliable: empirically,
    simply iterating `for s in prs.slides: for shp in s.shapes: ...` (the
    completely unrelated eyebrow-text lookup this script's own main()
    does) was enough to make prs.part.rels report the relationship set
    with 3 partnames duplicated on an otherwise byte-identical reload of
    the same file -- a python-pptx-internal lazy-loading/caching quirk.
    Reading the raw XML straight from the zip, once, before python-pptx's
    object model ever touches the relationship graph, sidesteps that
    category of bug entirely."""
    with zipfile.ZipFile(path) as z:
        try:
            data = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        except KeyError:
            return 0
    nums = [int(n) for n in re.findall(r'slide(\d+)\.xml"', data)]
    return max(nums) if nums else 0


def duplicate_bg_and_ids(prs, source_slide, keep_ids):
    """New slide on the same layout; copies the slide-level background and
    only the shapes whose shape_id is in keep_ids (used to lift just the
    eyebrow/title/intro/footer shapes off the PRODUCT template slide,
    leaving the hero-card and callout shapes behind for slides that don't
    want them)."""
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    dedupe_new_slide_partname(new_slide)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    source_bg = source_slide._element.find(f".//{NS_P}cSld/{NS_P}bg")
    if source_bg is not None:
        cSld = new_slide._element.find(f"{NS_P}cSld")
        spTree = cSld.find(f"{NS_P}spTree")
        cSld.insert(list(cSld).index(spTree), copy.deepcopy(source_bg))

    for shp in source_slide.shapes:
        if shp.shape_id in keep_ids:
            new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def set_text(slide, shape_id, text):
    for shp in slide.shapes:
        if shp.shape_id == shape_id and shp.has_text_frame:
            para = shp.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra._r.getparent().remove(extra._r)
            else:
                run = para.add_run()
                run.text = text
            return
    raise ValueError(f"shape_id {shape_id} not found")


def set_xfrm(shape_el, l, t, w, h):
    xfrm = shape_el.find(f".//{NS_A}xfrm")
    off = xfrm.find(f"{NS_A}off")
    ext = xfrm.find(f"{NS_A}ext")
    off.set("x", str(l))
    off.set("y", str(t))
    ext.set("cx", str(w))
    ext.set("cy", str(h))


def fit_hero(img_path):
    with Image.open(img_path) as im:
        w, h = im.size
    pic_w_avail = MAX_CARD_W - 2 * CARD_PAD
    pic_h_avail = MAX_CARD_H - 2 * CARD_PAD
    scale = min(pic_w_avail / w, pic_h_avail / h)
    pic_w, pic_h = round(w * scale), round(h * scale)
    card_w, card_h = pic_w + 2 * CARD_PAD, pic_h + 2 * CARD_PAD
    card_t = CARD_T + round((MAX_CARD_H - card_h) / 2)
    return {
        "card": (CARD_L, card_t, card_w, card_h),
        "pic": (CARD_L + CARD_PAD, card_t + CARD_PAD, pic_w, pic_h),
    }


def add_hero(slide, source_card_shape, img_path):
    geo = fit_hero(img_path)
    card_el = copy.deepcopy(source_card_shape._element)
    slide.shapes._spTree.append(card_el)
    l, t, w, h = geo["card"]
    set_xfrm(card_el, l, t, w, h)
    l, t, w, h = geo["pic"]
    slide.shapes.add_picture(img_path, Emu(l), Emu(t), Emu(w), Emu(h))
    return geo["card"]


def add_caption(slide, card_geo, text):
    l, t, w, h = card_geo
    box = slide.shapes.add_textbox(Emu(l), Emu(t + h + 60000), Emu(w), Emu(260000))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.name = "Calibri"
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def build_zoom_slide(prs, template_slide, source_card_shape, eyebrow, title, intro,
                      hero_path, caption, callouts):
    keep_ids = {2, 3, 4, 19, 20}
    slide = duplicate_bg_and_ids(prs, template_slide, keep_ids)
    set_text(slide, 2, eyebrow)
    set_text(slide, 3, title)
    set_text(slide, 4, intro)

    # Icon-row circle template (shape_id 7) and title/body text templates
    # (9/10) come from the PRODUCT slide too -- grab them once, then stamp
    # 3 copies at the fixed row positions, same as v7's approach.
    circle_src = next(s for s in template_slide.shapes if s.shape_id == 7)
    title_src = next(s for s in template_slide.shapes if s.shape_id == 9)
    body_src = next(s for s in template_slide.shapes if s.shape_id == 10)

    for i, (icon_path, ctitle, cbody) in enumerate(callouts):
        row_t = ICON_ROWS[i]
        circle_el = copy.deepcopy(circle_src._element)
        slide.shapes._spTree.append(circle_el)
        set_xfrm(circle_el, 8001000, row_t, 502920, 502920)
        slide.shapes.add_picture(icon_path, Emu(ICON_PIC_L), Emu(row_t + ICON_PIC_OFFSET_T), Emu(ICON_WH), Emu(ICON_WH))

        title_el = copy.deepcopy(title_src._element)
        slide.shapes._spTree.append(title_el)
        set_xfrm(title_el, TEXT_L, row_t - 27432, 2971800, 365760)
        title_el.find(f".//{NS_A}t").text = ctitle

        body_el = copy.deepcopy(body_src._element)
        slide.shapes._spTree.append(body_el)
        set_xfrm(body_el, TEXT_L, row_t + 204430, 2971800, 1234440)
        body_el.find(f".//{NS_A}t").text = cbody

    card_geo = add_hero(slide, source_card_shape, hero_path)
    add_caption(slide, card_geo, caption)
    return slide


# ---------------------------------------------------------------------------
# The 4-up SIGN / SEAL / QUOTE / DRAFT use-case summary slide
# ---------------------------------------------------------------------------

SLIDE_W = 12192000
MARGIN_L = 548640
N_COLS = 4
COL_GAP = 300000
COL_W = (SLIDE_W - 2 * MARGIN_L - (N_COLS - 1) * COL_GAP) // N_COLS

BADGE_WH = 700000
BADGE_T = 2450000
LABEL_T = BADGE_T + BADGE_WH + 200000
LABEL_H = 400000
BODY_T = LABEL_T + LABEL_H + 90000
BODY_H = 1900000

YELLOW = RGBColor(0xFD, 0xE0, 0x47)
NAVY_HEX = "0F172A"


def add_rounded_square(slide, l, t, wh, fill_rgb, radius_pct=0.28):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(wh), Emu(wh))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()
    try:
        shp.adjustments[0] = radius_pct
    except Exception:
        pass
    return shp


def add_text(slide, l, t, w, h, text, size_pt, color_rgb, bold=False, italic=False, align="l"):
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "c":
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = "Calibri"
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color_rgb
    return box


def build_usecases_slide(prs, template_slide):
    keep_ids = {2, 3, 4, 19, 20}
    slide = duplicate_bg_and_ids(prs, template_slide, keep_ids)
    # "USE CASES" rather than reusing "PRODUCT" (direct ask, 2026-08-08) --
    # the PRODUCT slide right before this one already uses that eyebrow, so
    # two in a row read like an accidental duplicate when flipping through.
    # This slide's actual content is literally the 4 CTA use cases, so the
    # label doubles as a truthful description.
    set_text(slide, 2, "USE CASES")
    set_text(slide, 3, "One product, four ways to get a document handled.")
    set_text(slide, 4, "Every path starts from the same New Document screen and ends in the same signer flow.")

    columns = [
        ("Signature", "Sign",
         "Upload a PDF you already have, place signature fields, and send it out — signed on any phone, no app required."),
        ("ShieldCheck", "Seal",
         "Seal a finished document to generate cryptographic proof — hash, timestamp, badge — that it's unaltered and identity-verified."),
        ("Receipt", "Quote",
         "Describe the job in plain language and get a line-item price quote, ready to review, edit, and send for signature."),
        ("Sparkles", "Draft",
         "Describe what you need in plain language and get a starting contract draft, ready to review and send."),
    ]

    for i, (icon_name, label, body) in enumerate(columns):
        col_l = MARGIN_L + i * (COL_W + COL_GAP)
        add_rounded_square(slide, col_l, BADGE_T, BADGE_WH, YELLOW)
        icon_wh = 340000
        icon_l = col_l + (BADGE_WH - icon_wh) // 2
        icon_t = BADGE_T + (BADGE_WH - icon_wh) // 2
        slide.shapes.add_picture(f"{APP_PUBLIC}/deck-icons/{icon_name}.png", Emu(icon_l), Emu(icon_t), Emu(icon_wh), Emu(icon_wh))
        add_text(slide, col_l, LABEL_T, COL_W, LABEL_H, label, 20, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        add_text(slide, col_l, BODY_T, COL_W, BODY_H, body, 12.5, RGBColor(0x94, 0xA3, 0xB8))

    return slide


def move_slides_after(prs, slides_to_move, after_index):
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    anchor = sld_ids[after_index]
    target_sldIds = []
    for slide in slides_to_move:
        for sldId in xml_slides:
            if int(sldId.get("id")) == slide.slide_id:
                target_sldIds.append(sldId)
                break
    for sldId in target_sldIds:
        xml_slides.remove(sldId)
    anchor_pos = list(xml_slides).index(anchor)
    for offset, sldId in enumerate(target_sldIds):
        xml_slides.insert(anchor_pos + 1 + offset, sldId)


def remove_slide(prs, slide_id):
    """Removes the <p:sldId> entry AND drops the underlying relationship
    in presentation.xml.rels, so the old slide's part is fully
    unreachable, not just undisplayed.

    HISTORY (2026-08-08, several rounds): the first version of this
    function deliberately skipped drop_rel(), reasoning that add_slide()'s
    naive len(sldIdLst)+1 partname scheme would otherwise collide with a
    still-displayed slide's part. That partname collision is now handled
    properly (dedupe_new_slide_partname() below assigns from a counter
    seeded via a raw XML scan, not by asking python-pptx what's "free").
    With that fixed, it turned out *leaving the old relationship in
    place* was its own, worse bug: a deck with any orphaned-but-still-
    related slide relationship reliably corrupted on the very next
    load-access-save cycle -- simply doing `list(prs.slides)[0]` (touching
    ANY slide, not even the orphaned ones) was enough to make python-pptx
    duplicate the orphaned slides' zip entries on save. Root cause not
    fully diagnosed (something in how the OPC loader resolves relationships
    once a target has no path to it via sldIdLst), but the fix is simple:
    don't leave any orphaned relationships lying around in the first
    place -- drop_rel() to fully sever the old slide, both display and
    relationship."""
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        if int(sldId.get("id")) == slide_id:
            rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            xml_slides.remove(sldId)
            prs.part.drop_rel(rId)
            return


def dedupe_new_slide_partname(new_slide):
    """Belt-and-suspenders around the len(sldIdLst)+1 partname bug
    described above: if add_slide() just handed back a partname that
    collides with an existing (different) part, rename it to a genuinely
    free /ppt/slides/slideN.xml before anything gets written out.

    UPDATED 2026-08-08 (twice): first compared against
    prs.part.package.iter_parts(), then against prs.part.rels -- both
    turned out to be unreliable (see max_slide_number_in_file's
    docstring above for specifics). Now just pulls the next number off a
    plain counter seeded once from a raw XML scan, no python-pptx
    object-graph involved."""
    from pptx.opc.packuri import PackURI

    _next_slide_n[0] += 1
    new_slide.part.partname = PackURI(f"/ppt/slides/slide{_next_slide_n[0]}.xml")


def main():
    _next_slide_n[0] = max_slide_number_in_file(SRC)
    prs = Presentation(SRC)
    slides_by_eyebrow = {}
    for s in prs.slides:
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t in ("PRODUCT", "SIGN", "SEAL", "DRAFT"):
                    slides_by_eyebrow.setdefault(t, s)
                    break

    template_slide = slides_by_eyebrow["PRODUCT"]
    old_sign = slides_by_eyebrow["SIGN"]
    old_seal = slides_by_eyebrow["SEAL"]
    old_draft = slides_by_eyebrow["DRAFT"]

    source_card_shape = next(s for s in template_slide.shapes if s.shape_id == 5)

    # Capture slide_ids up front, then remove the old standalone
    # SIGN/SEAL/DRAFT slides -- slide_id lookups scan the live sldIdLst,
    # so grabbing all three before mutating it avoids any risk of a
    # stale/invalidated reference partway through the loop.
    old_ids = [old_sign.slide_id, old_seal.slide_id, old_draft.slide_id]
    for sid in old_ids:
        remove_slide(prs, sid)

    usecases_slide = build_usecases_slide(prs, template_slide)

    sign_slide = build_zoom_slide(
        prs, template_slide, source_card_shape,
        eyebrow="SIGN",
        title="Rental signing that fits in the time it takes to hand over the keys.",
        intro="No app for the renter to download — scan a QR code, sign on their own phone, and drive off in minutes.",
        hero_path=f"{APP_PUBLIC}/hero-boat-jet-ski-rental-qr-signing.png",
        caption="Illustrative mockup of the shipped QR-to-sign flow — not a screenshot.",
        callouts=[
            (f"{APP_PUBLIC}/deck-icons/QrCode.png", "Scan to sign, no app",
             "The renter scans a QR code with their phone camera and signs directly in the browser — nothing to install."),
            (f"{APP_PUBLIC}/deck-icons/Clock.png", "Signed before they leave the counter",
             "The whole flow takes under a minute end to end, so operators aren't chasing signatures after the fact."),
            (f"{APP_PUBLIC}/deck-icons/ListChecks.png", "Every field already in place",
             "Fixed-fee and per-day rental terms are mapped to the right fields ahead of time — nothing to configure per rental."),
        ],
    )

    seal_slide = build_zoom_slide(
        prs, template_slide, source_card_shape,
        eyebrow="SEAL",
        title="Proof an invoice hasn't been touched since it was sent.",
        intro="A verification badge — hash, trusted timestamp, and a public verify page — embedded right on the document.",
        hero_path=f"{APP_PUBLIC}/hero-verified-badge-invoice.png",
        caption="Illustrative mockup of the shipped Seal / verify flow — not a screenshot.",
        callouts=[
            (f"{APP_PUBLIC}/deck-icons/Stamp.png", "Sealed, not just signed",
             "A cryptographic hash and an RFC 3161 trusted timestamp anchor the exact bytes of the invoice the moment it's sealed."),
            (f"{APP_PUBLIC}/deck-icons/ScanLine.png", "Anyone can verify it",
             "Scanning the badge opens a public page confirming the document is unchanged since it was sealed — no account needed."),
            (f"{APP_PUBLIC}/deck-icons/ShieldAlert.png", "Built to stop invoice fraud",
             "Protects freelancers and agencies against the classic scam: a client quietly editing bank details on a copy of a real invoice."),
        ],
    )

    product_index = list(prs.slides).index(template_slide)
    move_slides_after(prs, [usecases_slide, sign_slide, seal_slide], after_index=product_index)

    slides_in_order = list(prs.slides)
    for idx, slide in enumerate(slides_in_order, start=1):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip().isdigit():
                shp.text_frame.paragraphs[0].runs[0].text = str(idx)

    prs.save(DST)
    print("saved", DST, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
